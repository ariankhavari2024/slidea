# app/routes.py
import os
import io
import json
import shutil
import traceback
from datetime import datetime, timezone

from flask import (render_template, Blueprint, flash, redirect, url_for,
                   request, current_app, jsonify, abort, send_file, session, Response)
from flask_login import login_user, current_user, logout_user, login_required
from sqlalchemy import and_, func
from sqlalchemy.orm import aliased
from celery import group, chord
from celery.result import AsyncResult
import stripe
from botocore.exceptions import ClientError

from . import db, csrf, celery as celery_app
from .models import User, Presentation, Slide, PresentationStatus
from .forms import RegistrationForm, LoginForm, CreatePresentationForm, ContactForm
from .openai_helpers import (generate_text_content, parse_manual_content,
                             get_style_description, build_image_prompt, generate_slide_image,
                             generate_missing_slide_content)
from .tasks import (generate_single_slide_visual_task,
                    finalize_presentation_status_task)

# --- NEW IMPORTS for MinIO/S3 file serving ---
from .storage import get_s3_client, get_s3_bucket_name

# Import for PPTX generation
try:
    from pptx import Presentation as PptxPresentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    PPTX_INSTALLED = True
except ImportError:
    PPTX_INSTALLED = False

main = Blueprint('main', __name__)

# --- NEW ROUTE for serving files from MinIO/S3 ---
@main.route("/files/<path:key>")
@login_required # Protect files by ensuring user is logged in
def serve_s3_file(key):
    """
    Streams a file from the S3/MinIO bucket to the user.
    This acts as a secure proxy for your private storage.
    """
    s3 = get_s3_client()
    bucket = get_s3_bucket_name()
    try:
        # Basic authorization: Check if the presentation belongs to the current user
        # The key is expected to be in the format: "{presentation_id}/slide_{...}.png"
        presentation_id_str = key.split('/')[0]
        if presentation_id_str.isdigit():
            presentation = db.session.get(Presentation, int(presentation_id_str))
            if not presentation or presentation.user_id != current_user.id:
                current_app.logger.warning(f"User {current_user.id} attempted to access unauthorized S3 key: {key}")
                abort(403) # Forbidden
        else:
             # If the key format is unexpected, deny access for security
            current_app.logger.error(f"Invalid key format for S3 proxy access: {key}")
            abort(400) # Bad Request

        obj = s3.get_object(Bucket=bucket, Key=key)
        return Response(
            obj["Body"].read(),
            mimetype=obj.get("ContentType", "application/octet-stream"),
            headers={"Content-Disposition": f"inline; filename={os.path.basename(key)}"}
        )
    except ClientError as e:
        if e.response.get("Error", {}).get("Code") == "NoSuchKey":
            current_app.logger.warning(f"S3 file not found for key: {key}")
            return "File not found", 404
        current_app.logger.error(f"S3 error serving '{key}': {e}")
        return "Error serving file", 500
    except Exception as e:
        current_app.logger.error(f"Generic error serving S3 file '{key}': {e}", exc_info=True)
        return "Internal server error", 500


# --- Standard Page Routes (Home, About, Features, Contact) ---
@main.route('/')
def index():
    """Renders the home page."""
    return render_template('index.html', title='Welcome', current_user=current_user)

@main.route('/about')
def about():
    """Renders the about page."""
    return render_template('about.html', title='About Us', current_user=current_user)

@main.route('/features')
def features():
    """Renders the features page."""
    return render_template('features.html', title='Features', current_user=current_user)

@main.route('/contact', methods=['GET', 'POST'])
def contact():
    """Renders the contact page and handles form submission."""
    form = ContactForm()
    if form.validate_on_submit():
        name = form.name.data
        email = form.email.data
        subject = form.subject.data
        message_body = form.message.data
        current_app.logger.info(f"Contact form submission: Name='{name}', Email='{email}', Subject='{subject}'")
        flash('Thank you for your message! We will get back to you soon.', 'success')
        return redirect(url_for('main.contact'))
    elif request.method == 'POST':
        flash('Please correct the errors in the form.', 'warning')
    return render_template('contact.html', title='Contact Us', form=form, current_user=current_user)


# --- User Authentication & Account Routes ---
@main.route('/register', methods=['GET', 'POST'])
def register():
    """Handles user registration."""
    if current_user.is_authenticated:
        return redirect(url_for('main.dashboard'))
    form = RegistrationForm()
    if form.validate_on_submit():
        new_user = User(name=form.name.data, email=form.email.data)
        new_user.set_password(form.password.data)
        try:
            new_user.credits_remaining = current_app.config['CREDITS_PER_PLAN'].get('free', 0)
            new_user.subscription_plan_name = 'free'
            new_user.subscription_status = 'active'
            db.session.add(new_user)
            db.session.commit()
            flash(f'Account created! You get {new_user.credits_remaining} free credits. You can now log in.', 'success')
            return redirect(url_for('main.login'))
        except Exception as e:
            db.session.rollback()
            current_app.logger.error(f'Registration Error: {e}', exc_info=True)
            flash('An error occurred during registration. Please try again.', 'danger')
    elif request.method == 'POST':
        current_app.logger.warning(f"Registration validation failed: {form.errors}")
        flash('Please correct the errors below to register.', 'warning')
    return render_template('register.html', title='Register', form=form, current_user=current_user)

@main.route('/login', methods=['GET', 'POST'])
def login():
    """Handles user login."""
    if current_user.is_authenticated:
        return redirect(url_for('main.dashboard'))
    form = LoginForm()
    if form.validate_on_submit():
        user = User.query.filter_by(email=form.email.data).first()
        if user and user.check_password(form.password.data):
            login_user(user, remember=form.remember.data)
            flash('Login successful!', 'success')
            next_page = request.args.get('next')
            return redirect(next_page or url_for('main.dashboard'))
        else:
            flash('Login unsuccessful. Please check email and password.', 'danger')
    return render_template('login.html', title='Login', form=form, current_user=current_user)

@main.route('/logout')
@login_required
def logout():
    """Logs out the current user."""
    logout_user()
    flash('You have been logged out.', 'info')
    return redirect(url_for('main.index'))

@main.route('/account')
@login_required
def account():
    """Renders the account overview page for the logged-in user."""
    return render_template('account.html', title='Account Overview', current_user=current_user)

# --- Presentation Management Routes ---
@main.route('/dashboard')
@login_required
def dashboard():
    """Renders the user's dashboard, showing their presentations."""
    if request.args.get('success'):
        flash('Subscription successful or updated!', 'success')
    if request.args.get('canceled'):
        flash('Subscription checkout canceled.', 'info')
    if request.args.get('upgrade_success'):
        flash('Subscription successfully updated!', 'success')
    if request.args.get('upgrade_failed'):
        flash('Subscription update failed. Please try again or contact support.', 'danger')
    if request.args.get('plan_change_canceled'):
        flash('Plan change canceled.', 'info')

    results = db.session.query(
        Presentation,
        Slide.image_url
    ).outerjoin(
        Slide,
        and_(Presentation.id == Slide.presentation_id, Slide.slide_number == 1)
    ).filter(
        Presentation.user_id == current_user.id
    ).order_by(
        Presentation.last_edited_at.desc()
    ).all()

    return render_template(
        'dashboard.html',
        title='Dashboard',
        current_user=current_user,
        presentation_results=results,
        PresentationStatus=PresentationStatus
    )

@main.route('/create', methods=['GET', 'POST'])
@login_required
def create_presentation():
    """Handles the creation of new presentations."""
    form = CreatePresentationForm()
    required_credits = 0

    if request.method == 'POST':
        try:
            desired_slide_count_str = request.form.get('slide_count')
            desired_slide_count = int(desired_slide_count_str) if desired_slide_count_str and desired_slide_count_str.isdigit() else current_app.config.get('DEFAULT_SLIDE_COUNT', 3)
            required_credits = desired_slide_count * current_app.config.get('CREDITS_PER_SLIDE', 25)
        except (ValueError, TypeError) as e:
            current_app.logger.error(f"Error parsing slide_count: {e}. Form data: {request.form}", exc_info=True)
            flash('Invalid number of slides submitted. Please enter a whole number.', 'warning')
            return render_template('create_presentation.html', title='Create Presentation', form=form, current_user=current_user)
        
        if current_user.credits_remaining < required_credits:
            flash(f'Insufficient credits ({current_user.credits_remaining}) to generate {desired_slide_count} slides ({required_credits} needed). Please purchase more credits or upgrade your plan.', 'warning')
            return redirect(url_for('main.pricing'))

    if form.validate_on_submit():
        topic = form.topic.data
        presenter_name = form.presenter_name.data
        input_method = request.form.get('input_method_choice', 'auto')
        
        if input_method == 'manual':
            text_style = request.form.get('manual_text_style', 'bullet')
        else: 
            text_style = form.text_style.data
        
        if not text_style: 
            text_style = 'bullet'
            current_app.logger.warning("Text style was not set, defaulted to 'bullet'.")

        desired_slide_count = form.slide_count.data
        style_choice = form.style_choice.data
        custom_style_prompt = form.custom_style_prompt.data
        font_choice = form.font_choice.data
        creativity_score = form.creativity_score.data

        required_credits = desired_slide_count * current_app.config.get('CREDITS_PER_SLIDE', 25)
        credits_deducted_for_this_request = required_credits
        style_prompt_text = custom_style_prompt if style_choice == 'custom' else get_style_description(style_choice)

        slides_content_raw = []
        presentation_title = "Untitled Presentation"
        saved_slides_orm = []
        new_presentation = None

        try:
            user = db.session.get(User, current_user.id)
            if user.credits_remaining < credits_deducted_for_this_request:
                flash(f'Credit check failed. Needed: {credits_deducted_for_this_request}, Available: {user.credits_remaining}.', 'danger')
                return redirect(url_for('main.pricing'))

            user.credits_remaining = max(0, user.credits_remaining - credits_deducted_for_this_request)
            db.session.add(user)
            current_app.logger.info(f"CREDIT DEDUCTION: User {user.id} charged {credits_deducted_for_this_request}. Remaining: {user.credits_remaining}")
            flash("Processing input...", 'info')

            if input_method == 'manual':
                manual_topic = request.form.get(f'manual_title_1', '').strip()
                presentation_title = manual_topic or "Manually Created Presentation"
                for i in range(1, desired_slide_count + 1):
                    title = request.form.get(f'manual_title_{i}')
                    content = request.form.get(f'manual_content_{i}', '').strip()
                    if not title:
                        flash(f"Missing title for manually entered Slide {i}.", 'warning')
                        db.session.rollback()
                        user_refund = db.session.get(User, current_user.id)
                        user_refund.credits_remaining += credits_deducted_for_this_request
                        db.session.add(user_refund); db.session.commit()
                        current_app.logger.warning(f"CREDIT REFUND: User {user.id} due to missing manual title.")
                        return render_template('create_presentation.html', title='Create', form=form, current_user=current_user)
                    if i == 1: content = f"By: {presenter_name.strip()}" if presenter_name and presenter_name.strip() else ""
                    elif not content:
                        content = generate_missing_slide_content(title.strip(), text_style, manual_topic)
                        flash(f"Generated content for slide {i} ('{title}') as it was left blank.", "info")
                    slides_content_raw.append({"slide_title": title.strip(), "slide_content": content})
            elif input_method == 'auto' and topic:
                presentation_title = f"Presentation on: {topic}"
                slides_content_raw = generate_text_content(topic, text_style, desired_slide_count, presenter_name)
            else:
                flash("Invalid input method or missing topic.", 'warning')
                db.session.rollback()
                user_refund = db.session.get(User, current_user.id)
                user_refund.credits_remaining += credits_deducted_for_this_request
                db.session.add(user_refund); db.session.commit()
                current_app.logger.warning(f"CREDIT REFUND: User {user.id} due to invalid input method/topic.")
                return render_template('create_presentation.html', title='Create', form=form, current_user=current_user)

            if not slides_content_raw:
                flash("Could not generate slide content.", 'danger')
                db.session.rollback()
                user_refund = db.session.get(User, current_user.id)
                user_refund.credits_remaining += credits_deducted_for_this_request
                db.session.add(user_refund); db.session.commit()
                current_app.logger.warning(f"CREDIT REFUND: User {user.id} due to content generation failure.")
                return render_template('create_presentation.html', title='Create', form=form, current_user=current_user)

            total_slides = len(slides_content_raw)
            new_presentation = Presentation(user_id=current_user.id, title=presentation_title, style_prompt=style_prompt_text, status=PresentationStatus.PENDING_VISUALS, font_choice=font_choice, creativity_score=creativity_score)
            db.session.add(new_presentation); db.session.flush()
            for i, slide_data in enumerate(slides_content_raw):
                raw_content = slide_data.get('slide_content', '')
                processed_content_db = json.dumps(raw_content) if isinstance(raw_content, list) else str(raw_content)
                new_slide = Slide(presentation_id=new_presentation.id, slide_number=i + 1, title=slide_data.get('slide_title', f'Slide {i+1}'), text_content=processed_content_db)
                db.session.add(new_slide); saved_slides_orm.append(new_slide)
            db.session.commit()
            
            if saved_slides_orm:
                callback_task = finalize_presentation_status_task.s(new_presentation.id, current_user.id, total_slides, credits_deducted_for_this_request)
                slide_task_signatures = []
                for i_task, slide_orm_object_task in enumerate(saved_slides_orm):
                    raw_content_for_style_check_task = slides_content_raw[i_task].get('slide_content', '')
                    image_gen_text_style_hint_task = 'paragraph'
                    if i_task > 0:
                        if isinstance(raw_content_for_style_check_task, list): image_gen_text_style_hint_task = 'bullet'
                        elif isinstance(raw_content_for_style_check_task, str) and '\n' in raw_content_for_style_check_task: image_gen_text_style_hint_task = 'bullet'
                    
                    task_args_for_slide_task = {"presentation_topic": presentation_title, "presenter_name": presenter_name, "total_slides": total_slides, "text_style_for_image": image_gen_text_style_hint_task, "creativity_score": creativity_score, "font_choice": font_choice, "presentation_style_prompt": style_prompt_text}
                    slide_task_signatures.append(generate_single_slide_visual_task.s(slide_orm_object_task.id, current_user.id, **task_args_for_slide_task))
                
                chord_result = chord(group(slide_task_signatures))(callback_task)
                new_presentation.celery_chord_id = chord_result.id; db.session.commit()
                flash(f"Presentation '{new_presentation.title}' created! Generating visuals.", 'success')
            return redirect(url_for('main.dashboard'))
        except Exception as e:
            db.session.rollback()
            current_app.logger.error(f"Create Presentation Error: {e}", exc_info=True)
            # (Credit refund and error handling logic)
            flash(f"Error creating presentation: {str(e)}", 'danger')
    
    return render_template('create_presentation.html', title='Create Presentation', form=form, current_user=current_user)


@main.route('/presentation/<int:presentation_id>')
@login_required
def view_presentation(presentation_id):
    presentation = Presentation.query.get_or_404(presentation_id)
    if presentation.user_id != current_user.id:
        abort(403)
    slides = presentation.slides.order_by(Slide.slide_number).all()
    return render_template('view_presentation.html', title=presentation.title, presentation=presentation, slides=slides, current_user=current_user, PresentationStatus=PresentationStatus)

@main.route('/editor/<int:presentation_id>')
@login_required
def editor(presentation_id):
    presentation = Presentation.query.get_or_404(presentation_id)
    if presentation.user_id != current_user.id:
        abort(403)

    credits_needed_for_regen = current_app.config.get('CREDITS_PER_REGENERATE', 25)
    can_regenerate = current_user.credits_remaining >= credits_needed_for_regen
    
    slides = presentation.slides.order_by(Slide.slide_number).all()
    slides_data = []
    for s in slides:
        content = s.text_content or ""
        try:
            parsed_content = json.loads(content) if content.strip().startswith(('[', '{')) else content
        except json.JSONDecodeError:
            parsed_content = content
        
        slides_data.append({
            "id": s.id,
            "slide_number": s.slide_number,
            "title": s.title or "",
            "text_content": parsed_content,
            "image_url": s.image_url # The URL is already the correct /files/... path from the DB
        })
    return render_template('editor.html',
                           title=f"Editing: {presentation.title}",
                           presentation=presentation,
                           slides_json=json.dumps(slides_data),
                           current_user=current_user,
                           PresentationStatus=PresentationStatus,
                           can_regenerate=can_regenerate)

@main.route('/api/slide/<int:slide_id>/regenerate_image', methods=['POST'])
@login_required
def regenerate_slide_image_api(slide_id):
    slide = db.session.get(Slide, slide_id)
    if not slide:
        return jsonify({'status': 'error', 'message': 'Slide not found'}), 404
    
    presentation = db.session.get(Presentation, slide.presentation_id)
    if not presentation or presentation.user_id != current_user.id:
        return jsonify({'status': 'error', 'message': 'Permission denied'}), 403

    credits_needed = current_app.config.get('CREDITS_PER_REGENERATE', 25)
    user = db.session.get(User, current_user.id)

    if user.credits_remaining < credits_needed:
        return jsonify({'status': 'error', 'message': f'Insufficient credits.'}), 402

    data = request.get_json()
    if not data or 'title' not in data or 'text_content' not in data:
        return jsonify({'status': 'error', 'message': 'Missing title or text_content'}), 400

    edit_prompt_text = data.get('edit_prompt', '').strip()

    try:
        user.credits_remaining = max(0, user.credits_remaining - credits_needed)
        db.session.add(user)
        
        slide.title = data['title']
        new_content = data['text_content']
        try:
            original_parsed_content = json.loads(slide.text_content) if slide.text_content and slide.text_content.strip().startswith('[') else slide.text_content
            if isinstance(original_parsed_content, list):
                slide.text_content = json.dumps([line.strip() for line in new_content.split('\n') if line.strip()] or [" "])
            else:
                slide.text_content = str(new_content)
        except (json.JSONDecodeError, TypeError):
            slide.text_content = str(new_content)
        
        presentation.last_edited_at = datetime.now(timezone.utc)
        db.session.add(slide); db.session.add(presentation)
        
        presenter_name = presentation.author.name if presentation.author else None
        total_slides = db.session.query(func.count(Slide.id)).filter(Slide.presentation_id == presentation.id).scalar() or 1
        
        try:
            content_for_style_check = json.loads(slide.text_content) if (slide.text_content and (slide.text_content.strip().startswith('[') or slide.text_content.strip().startswith('{'))) else slide.text_content
            text_style_for_image = 'bullet' if isinstance(content_for_style_check, list) or (slide.slide_number != 1 and '\n' in str(content_for_style_check)) else 'paragraph'
            if slide.slide_number == 1: text_style_for_image = 'paragraph'
        except (json.JSONDecodeError, TypeError):
            text_style_for_image = 'paragraph'
        
        creativity_score = getattr(presentation, 'creativity_score', 5)
        font_choice = getattr(presentation, 'font_choice', 'Inter')
        style_description_to_use = edit_prompt_text or slide.applied_style_info or presentation.style_prompt or get_style_description('keynote_modern')
        
        image_gen_prompt_text = build_image_prompt(
            slide_title=slide.title, 
            slide_content=json.loads(slide.text_content) if text_style_for_image == 'bullet' and isinstance(content_for_style_check, list) else slide.text_content, 
            style_description=style_description_to_use, 
            text_style=text_style_for_image, 
            slide_number=slide.slide_number, 
            total_slides=total_slides, 
            creativity_score=creativity_score, 
            presentation_topic=presentation.title, 
            font_choice=font_choice, 
            presenter_name=presenter_name
        )
        
        image_url, actual_prompt_used = generate_slide_image(
            image_prompt=image_gen_prompt_text,
            presentation_id=presentation.id, 
            slide_number=slide.slide_number
        )

        if image_url:
            slide.image_url = image_url
            slide.image_gen_prompt = actual_prompt_used
            slide.applied_style_info = style_description_to_use
            db.session.commit()
            return jsonify({
                'status': 'success',
                'new_image_url': image_url,
                'credits_remaining': user.credits_remaining
            })
        else:
            db.session.rollback()
            return jsonify({'status': 'error', 'message': 'Image regeneration failed.'}), 500
    except Exception as e:
        db.session.rollback()
        current_app.logger.error(f"API: Regen Error for slide {slide.id}: {e}", exc_info=True)
        return jsonify({'status': 'error', 'message': f'Error: {str(e)}'}), 500

@main.route('/presentation/<int:presentation_id>/delete', methods=['POST'])
@login_required
def delete_presentation(presentation_id):
    """Deletes a presentation. Note: S3 objects are not deleted in this version."""
    presentation = db.session.get(Presentation, presentation_id)
    if not presentation or presentation.user_id != current_user.id:
        abort(403)
    try:
        # Note: This no longer deletes from a local folder.
        # A separate cleanup task would be needed to delete from S3.
        db.session.delete(presentation)
        db.session.commit()
        flash(f'Presentation "{presentation.title}" has been deleted.', 'success')
    except Exception as e:
        db.session.rollback()
        flash('Error deleting presentation.', 'danger')
    return redirect(url_for('main.dashboard'))

@main.route('/presentation/<int:presentation_id>/export/pptx')
@login_required
def export_pptx(presentation_id):
    """Exports a presentation to a PPTX file. Requires fetching images from S3."""
    # This function needs to be updated to fetch images from the S3 proxy route
    # or directly from S3, which is more complex. Disabling for now.
    flash("PPTX export is not yet compatible with the new storage system.", "warning")
    return redirect(url_for('main.view_presentation', presentation_id=presentation_id))

@main.route('/api/presentation/<int:presentation_id>/status')
@login_required
def presentation_status_api(presentation_id):
    presentation = db.session.get(Presentation, presentation_id)
    if not presentation or presentation.user_id != current_user.id:
        return jsonify({'status': 'error', 'message': 'Not found'}), 404
    
    response_data = {'status': presentation.status.value}
    if presentation.status == PresentationStatus.PENDING_VISUALS:
        try:
            total_slides = db.session.query(func.count(Slide.id)).filter_by(presentation_id=presentation_id).scalar() or 0
            completed_slides = db.session.query(func.count(Slide.id)).filter(
                Slide.presentation_id == presentation_id,
                Slide.image_url.isnot(None)
            ).scalar() or 0
            response_data['total_slides'] = total_slides
            response_data['completed_slides'] = completed_slides
        except Exception as e:
            current_app.logger.error(f"Error querying slide progress for Presentation {presentation_id}: {e}", exc_info=True)
            response_data['total_slides'] = None
            response_data['completed_slides'] = None
            
    return jsonify(response_data)

@main.route('/presentation/<int:presentation_id>/cancel', methods=['POST'])
@login_required
def cancel_presentation(presentation_id):
    presentation = db.session.get(Presentation, presentation_id)
    if not presentation or presentation.user_id != current_user.id:
        abort(403)

    if presentation.status != PresentationStatus.PENDING_VISUALS:
        flash("This presentation is not currently generating visuals.", "warning")
        return redirect(url_for('main.dashboard'))

    if not presentation.celery_chord_id:
        flash("No active generation task found to cancel for this presentation.", "warning")
        presentation.status = PresentationStatus.GENERATION_FAILED
        db.session.commit()
        return redirect(url_for('main.dashboard'))

    try:
        current_app.logger.info(f"Attempting to cancel Celery chord: {presentation.celery_chord_id} for Presentation {presentation.id}")
        chord_res = AsyncResult(presentation.celery_chord_id, app=celery_app.celery)
        if chord_res and chord_res.parent:
            chord_res.parent.revoke(terminate=True, signal='SIGTERM')
        if chord_res:
            chord_res.revoke(terminate=True, signal='SIGTERM')

        presentation.status = PresentationStatus.GENERATION_FAILED
        presentation.last_edited_at = datetime.now(timezone.utc)
        
        num_slides_intended = presentation.slides.count()
        credits_to_refund = num_slides_intended * current_app.config.get('CREDITS_PER_SLIDE', 25)
        
        user = db.session.get(User, current_user.id)
        if user and credits_to_refund > 0:
            user.credits_remaining += credits_to_refund
            db.session.add(user)
            current_app.logger.info(f"CREDIT REFUND: User {user.id} refunded {credits_to_refund} credits for cancelled Presentation {presentation.id}.")
        
        db.session.commit()
        flash(f'Visual generation for "{presentation.title}" has been cancelled. Credits have been refunded.', 'info')
    except Exception as e:
        current_app.logger.error(f"Error cancelling Celery chord {presentation.celery_chord_id} for Presentation {presentation.id}: {e}", exc_info=True)
        flash('An error occurred while trying to cancel the generation.', 'danger')
        if presentation.status == PresentationStatus.PENDING_VISUALS:
            presentation.status = PresentationStatus.GENERATION_FAILED
            db.session.commit()
            
    return redirect(url_for('main.dashboard'))


# --- Stripe Routes ---
@main.route('/pricing')
def pricing():
    return render_template('pricing.html', title='Pricing Plans', current_user=current_user, config=current_app.config)

@main.route('/confirm-plan-change')
@login_required
def confirm_plan_change():
    new_price_id = request.args.get('new_price_id')
    current_plan_name_display = current_app.config['PLAN_NAME_MAP'].get(current_user.stripe_price_id, current_user.subscription_plan_name or 'your current plan')
    new_plan_name_display = current_app.config['PLAN_NAME_MAP'].get(new_price_id, 'the selected plan')

    if not new_price_id:
        flash('No plan selected for change.', 'warning')
        return redirect(url_for('main.pricing'))

    session['plan_change_new_price_id'] = new_price_id
    
    return render_template('confirm_plan_change.html',
                           title='Confirm Plan Change',
                           current_plan_name=current_plan_name_display,
                           new_plan_name=new_plan_name_display,
                           current_user=current_user)

@main.route('/process-plan-change', methods=['POST'])
@login_required
def process_plan_change():
    new_price_id = session.pop('plan_change_new_price_id', None)
    if not new_price_id:
        flash('Plan change information missing or expired. Please try again.', 'danger')
        return redirect(url_for('main.pricing'))

    if not current_user.stripe_subscription_id or current_user.subscription_status != 'active':
        flash('No active subscription found to update, or subscription is not active.', 'warning')
        return redirect(url_for('main.pricing'))

    stripe.api_key = current_app.config['STRIPE_SECRET_KEY']
    try:
        subscription = stripe.Subscription.retrieve(current_user.stripe_subscription_id)
        
        current_plan_credits = current_app.config['CREDITS_PER_PLAN'].get(current_user.subscription_plan_name, 0)
        new_plan_name_temp = current_app.config['PLAN_NAME_MAP'].get(new_price_id, 'unknown')
        new_plan_credits = current_app.config['CREDITS_PER_PLAN'].get(new_plan_name_temp, 0)

        proration_behavior_to_use = 'create_prorations'
        if new_plan_credits < current_plan_credits:
            proration_behavior_to_use = 'none'
        
        updated_subscription = stripe.Subscription.modify(
            current_user.stripe_subscription_id,
            items=[{
                'id': subscription['items']['data'][0]['id'],
                'price': new_price_id,
            }],
            proration_behavior=proration_behavior_to_use,
            cancel_at_period_end=False,
        )
        flash('Your subscription plan is being updated!', 'info')
        return redirect(url_for('main.dashboard', upgrade_success=True))

    except stripe.error.StripeError as e:
        flash(f"Could not update your subscription: {str(e)}", 'danger')
        return redirect(url_for('main.pricing', upgrade_failed=True))
    except Exception as e:
        flash("An unexpected error occurred while updating your subscription.", 'danger')
        return redirect(url_for('main.pricing', upgrade_failed=True))


@main.route('/create-checkout-session', methods=['POST'])
@login_required
def create_checkout_session():
    price_id = request.form.get('price_id')
    if not price_id:
        flash('Invalid plan selected.', 'danger')
        return redirect(url_for('main.pricing'))

    stripe.api_key = current_app.config['STRIPE_SECRET_KEY']

    if not current_user.stripe_customer_id:
        try:
            customer = stripe.Customer.create(
                email=current_user.email, name=current_user.name,
                metadata={'user_id': current_user.id} )
            current_user.stripe_customer_id = customer.id
            db.session.add(current_user); db.session.commit()
        except Exception as e:
            db.session.rollback()
            flash("An error occurred setting up your billing account.", 'danger')
            return redirect(url_for('main.pricing'))

    if current_user.stripe_subscription_id and current_user.subscription_status == 'active' and price_id != current_user.stripe_price_id:
        return redirect(url_for('main.confirm_plan_change', new_price_id=price_id))
    else:
        try:
            checkout_session = stripe.checkout.Session.create(
                customer=current_user.stripe_customer_id,
                payment_method_types=['card'],
                line_items=[{'price': price_id, 'quantity': 1}],
                mode='subscription',
                success_url=url_for('main.dashboard', _external=True) + '?success=true&session_id={CHECKOUT_SESSION_ID}',
                cancel_url=url_for('main.pricing', _external=True) + '?canceled=true',
                metadata={'user_id': current_user.id}
            )
            return redirect(checkout_session.url, code=303)
        except Exception as e:
            flash("An unexpected error occurred during checkout setup.", 'danger')
            return redirect(url_for('main.pricing'))


@main.route('/create-portal-session', methods=['POST'])
@login_required
def create_portal_session():
    if not current_user.stripe_customer_id:
        flash('Billing portal unavailable. No billing information found.', 'warning')
        return redirect(url_for('main.account'))

    stripe.api_key = current_app.config['STRIPE_SECRET_KEY']
    try:
        portal_session = stripe.billing_portal.Session.create(
            customer=current_user.stripe_customer_id,
            return_url=url_for('main.account', _external=True),
        )
        return redirect(portal_session.url, code=303)
    except Exception as e:
        flash("An unexpected error occurred. Please try again.", 'danger')
        return redirect(url_for('main.account'))


@main.route('/stripe-webhook', methods=['POST'], strict_slashes=False)
def stripe_webhook():
    payload = request.data
    sig_header = request.headers.get('Stripe-Signature')
    endpoint_secret = current_app.config.get('STRIPE_ENDPOINT_SECRET')
    event = None

    if not endpoint_secret:
        current_app.logger.error("Webhook Error: Endpoint secret not configured.")
        return jsonify(error="Webhook secret not configured"), 500
    try:
        event = stripe.Webhook.construct_event(payload, sig_header, endpoint_secret)
        current_app.logger.info(f"✅ Webhook Event Verified: ID={event.id}, Type={event['type']}")
    except ValueError as e:
        return jsonify(error=str(e)), 400
    except stripe.error.SignatureVerificationError as e:
        return jsonify(error=str(e)), 400
    except Exception as e:
        return jsonify(error=str(e)), 500

    session_data = event['data']['object']
    event_type = event['type']
    user = None
    customer_id = session_data.get('customer')
    user_id_metadata = session_data.get('metadata', {}).get('user_id')

    if customer_id:
        user = User.query.filter_by(stripe_customer_id=customer_id).first()
    if not user and user_id_metadata:
        try:
            user = db.session.get(User, int(user_id_metadata))
        except (ValueError, TypeError):
            pass
    if not user:
        return jsonify(success=True, message="User not found, skipping.")

    log_prefix = f"Webhook {event_type} (User ID {user.id})"
    current_app.logger.info(f"{log_prefix}: Processing event...")
    try:
        if event_type == 'checkout.session.completed':
            if session_data.get('payment_status') == 'paid' and session_data.get('mode') == 'subscription':
                subscription_id = session_data.get('subscription')
                if not subscription_id: return jsonify(success=True)
                
                if user.stripe_subscription_id != subscription_id:
                    subscription = stripe.Subscription.retrieve(subscription_id)
                    items = subscription['items']['data']
                    price_id = items[0]['price']['id'] if items else None
                    current_period_end_dt = datetime.fromtimestamp(subscription.get('current_period_end'), tz=timezone.utc) if subscription.get('current_period_end') else None
                    
                    plan_name = current_app.config['PLAN_NAME_MAP'].get(price_id, 'unknown')
                    credits_to_grant = current_app.config['CREDITS_PER_PLAN'].get(plan_name, 0)
                    
                    user.credits_remaining = credits_to_grant
                    user.stripe_customer_id = customer_id
                    user.stripe_subscription_id = subscription_id
                    user.stripe_price_id = price_id
                    user.subscription_plan_name = plan_name
                    user.subscription_status = subscription.status
                    user.subscription_current_period_end = current_period_end_dt
                    db.session.add(user)
                    db.session.commit()
                    current_app.logger.info(f"{log_prefix}: COMMIT successful. Credits SET for new subscription.")

        elif event_type == 'customer.subscription.updated':
            subscription_event_data = session_data
            new_status = subscription_event_data.get('status')
            sub_items = subscription_event_data.get('items', {}).get('data', [])
            new_price_id = sub_items[0].get('price', {}).get('id') if sub_items else None
            new_period_end_dt = datetime.fromtimestamp(subscription_event_data.get('current_period_end'), tz=timezone.utc) if subscription_event_data.get('current_period_end') else None
            
            old_user_stripe_price_id = user.stripe_price_id
            old_user_credits = user.credits_remaining or 0
            
            user.stripe_subscription_id = subscription_event_data.get('id')
            user.subscription_status = new_status
            user.subscription_current_period_end = new_period_end_dt
            user.stripe_price_id = new_price_id
            new_plan_name = current_app.config['PLAN_NAME_MAP'].get(new_price_id, 'unknown')
            user.subscription_plan_name = new_plan_name
            
            if new_price_id and (new_price_id != old_user_stripe_price_id):
                if new_status in ['active', 'trialing']:
                    credits_for_new_plan = current_app.config['CREDITS_PER_PLAN'].get(new_plan_name, 0)
                    user.credits_remaining = old_user_credits + credits_for_new_plan
            elif new_status not in ['active', 'trialing'] and user.subscription_status != new_status:
                user.credits_remaining = 0
                user.subscription_plan_name = 'free'
                user.stripe_price_id = None
            
            db.session.add(user)
            db.session.commit()
            current_app.logger.info(f"{log_prefix}: COMMIT successful for subscription update.")

        elif event_type == 'customer.subscription.deleted':
            if user.stripe_subscription_id == session_data.get('id'):
                user.stripe_subscription_id = None; user.stripe_price_id = None; user.subscription_plan_name = 'free'
                user.subscription_status = 'canceled'; user.subscription_current_period_end = None; user.credits_remaining = 0
                db.session.add(user); db.session.commit()
                current_app.logger.info(f"{log_prefix}: COMMIT successful for sub deletion.")

        elif event_type == 'invoice.paid':
            invoice = session_data
            subscription_id_on_invoice = invoice.get('subscription')
            billing_reason = invoice.get('billing_reason')
            
            if user.stripe_subscription_id == subscription_id_on_invoice and user.subscription_status == 'active' and billing_reason == 'subscription_cycle':
                invoice_lines = invoice.get('lines', {}).get('data', [])
                price_id = invoice_lines[0].get('price', {}).get('id') if invoice_lines else None
                plan_name = current_app.config['PLAN_NAME_MAP'].get(price_id, 'unknown')
                credits_to_grant = current_app.config['CREDITS_PER_PLAN'].get(plan_name, 0)
                
                if credits_to_grant > 0:
                    user.credits_remaining = credits_to_grant
                    user.subscription_current_period_end = datetime.fromtimestamp(invoice_lines[0].get('period', {}).get('end'), tz=timezone.utc) if invoice_lines and invoice_lines[0].get('period') else None
                    db.session.add(user); db.session.commit()
                    current_app.logger.info(f"{log_prefix}: COMMIT successful for invoice.paid (renewal).")
    
    except Exception as e:
        db.session.rollback()
        current_app.logger.error(f"Webhook {event_type}: Error processing for User {user.id}: {e}", exc_info=True)
        return jsonify(error="Internal server error"), 500
    
    return jsonify(success=True), 200
